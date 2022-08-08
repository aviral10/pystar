# UI library
import streamlit as st
#for pthon3.10
import collections 
import collections.abc

# Pandas for dataframes
import os
import pandas as pd
import numpy as np
from math import isnan

# for pptx
from pptx import Presentation
from pptx.util import Inches, Cm, Pt

# importing all utility functions
from src.util import *
import streamlit.components.v1 as components  # Import Streamlit



# Page header
st.set_page_config(
    page_title="pyStAR", page_icon="🥂"
)

# Removing Streamlit watermark
hide_streamlit_style = """
            <style>
            #MainMenu {visibility: hidden;}
            footer {visibility: hidden;}
            </style>
            """
st.markdown(hide_streamlit_style, unsafe_allow_html=True) 

# logo and header
# <h1 style="font-size: 25px"> Excel to PPT Automation </h1>
components.html(f"""
    <div style="text-align: right;height: 100%;width:100%;font-family: Adobe Clean;
    margin-top: -20px;
    background-repeat: no-repeat;
    "> 
        <h1 style="font-size: 50px"> pyStAR 🌠 </h1>
        
    </div>
    <div>
        <hr style="width:100%;text-align:left;margin-left:0;color:black;background-color:black;height: 2px;border-radius: 25px;">
    </div>
""", height=170)

# st.markdown("# pyStAR ⭐")
# st.markdown("# Excel to PPT Automation")

#globals
default_template = 'data/template.pptx'
template_file = default_template
prs = Presentation(template_file)    # PPT Presentation object
prs.slide_width = Cm(33.858)    # Width of the slide
prs.slide_height = Cm(19.05)    # Height of the slide
sz = 6    # Global parameter(later overwritten) for setting rows per slide
limit = -1    # Dev mode only feature

# user input
user_start_location = None    # Starting cell address, e.g: A13
break_into_slides = "no"    # If we want to break data into opencases and closed cases
global_compare_false_val = '-'    # This value is considered as False value and rows corresponding to this value are put in closed case slide.
global_null_value = '-'    # This value acts as empty value in the dataframe as well as on the PPT slides
openCaseCol = None    # This column is used to segreagate opencases and closedcases
opcsheads = []    # Headers/Column names for opencases
cscsheads = []    # Headers/Column names for closedcases
opencases = []    # Rows to be displayed on open case slide
closedcases = []    # Rows to be displayed on close case slide
ext_tab_arr = []    # Array containing data to be displayed on extra slides
ext_cha_arr = []    # Array containing data from which charts are to be created
openCase_rows_per_slide = sz    # Number of rows to be displayed per slide in opencases
closeCase_rows_per_slide = sz    # Number of rows to be displayed per slide in closedcases

beginx = None    # Utility variable that stores staring row address of the data
beginy = None    # Utility variable that stores staring column address of the data
dataframe = None    # Will be updated with the Final Data after it is extracted from excel file

def get_data(file):
    """
    Extracts data from an Excel file into a dataframe.

    @param file: Path to the excel file
    @return: Pandas dataframe
    """
    # global references
    global beginx, beginy

    # reading excel file, extracting all data
    df = pd.read_excel(file, dtype=str)

    # handling cases where first row of the excel file may be empty, or the data is offset by some rows.
    beginy = max(0, beginy-1)
    if beginx > 2:
        beginx -= 2
    else:
        beginx -= 1

    # truncating dataframe according to the offsets: beginx, beginy
    df = df.iloc[beginx:,beginy:]

    # Evaluating headers(column titles)
    new_header = df.columns
    if beginx != 0:
        new_header = df.iloc[0] 
        df = df[1:] 
    df.columns = new_header

    # replace null values in the dataframe with global_null_value
    df = df.replace(np.nan, global_null_value)
    return df

def extract_rows(df, break_slides=True, headers=[]):
    """
    Extracts rows from the excel file. 
    The parameter break_slides determines whether opencase and closed case segregation is required.
    If break_slides is True, data extracted from @Dataframe df is divided into 2 @Array opencases and @Array closedcases.
    If break_sldies is False, data extracted from @Dataframe df is simply put into @Array entries.

    @param df: dataframe
    @param break_slides: boolean value to check if segregation of open cases and closed cases is required.
    @param headers: ONLY REQUIRED when break_slides is false, contains column titles for each slide
    @return: @Array Opencases and @Array Closedcases when break_slides is True, else returns @Array entries.
    """
    
    # check for the break_slides condition
    if break_slides == True:
        # initialize opencases and closed cases arrays
        opencases = []
        closedcases = []

        # IMPORTANT : if some column does not exist in the dataframe, 
        # it is created and filled with null values
        # checking for both headers
        for ele in opcsheads:
            if ele not in dataframe.columns:
                dataframe[ele] = global_null_value
        for ele in cscsheads:
            if ele not in dataframe.columns:
                dataframe[ele] = global_null_value
        
        # get the column on the basis of which data is segregated into open and closed cases
        open_cases_bool = dataframe[openCaseCol]

        opencaseheaders = opcsheads
        closedcaseheaders = cscsheads
        
        # extracting data according to the specified column header for opencases
        opendf = df[opencaseheaders]

        # first entry in opencases are the column titles, these titles make up the topmost row displayed in every slide.
        opencases.append(opendf.columns.tolist())
        
        # extract data for opencases
        for index, row in opendf.iterrows():
            temp = []
            # comparator_break is used to segregate the data
            if comparator_break(open_cases_bool[index], global_compare_false_val):
                continue
            # if the row is valid, this goes in @Array opencases 
            for ele in row:
                temp.append(str(ele))
            opencases.append(temp)
        
        # extracting data according to the specified column header for closedcases
        closedf = df[closedcaseheaders]

        # first entry in closedcases are the column titles, these titles make up the topmost row displayed in every slide.
        closedcases.append(closedf.columns.tolist())
        # extract data for closedcases
        for index, row in closedf.iterrows():
            temp = []
            # comparator_break is used to segregate the data
            if not comparator_break(open_cases_bool[index], global_compare_false_val):
                continue

            # if the row is valid, this goes in @Array closedcases
            for ele in row:
                temp.append(str(ele))
            closedcases.append(temp)
        
        # return the @Array opencases and @Array closedcases
        return opencases, closedcases
        
    else:
        # if break_sides is False
        # initialise @Array entries
        entries = []

        # if invalid headers were passed return empty entries
        if len(headers) == 0:
            return [[]]
        else:
            # IMPORTANT : if some column does not exist in the dataframe, 
            # it is created and filled with null values
            # checking for both headers
            for ele in headers:
                if ele not in dataframe.columns:
                    dataframe[ele] = global_null_value

        # extracting data according to the specified column headers
        opendf = df[headers]
        # first entry in entries are the column titles, these titles make up the topmost row displayed in every slide.
        entries.append(opendf.columns.tolist())
        for index, row in opendf.iterrows():
            temp = []
            for ele in row:
                temp.append(str(ele))
            entries.append(temp)
        
        return entries

def present_on_slide(data,sz=6, title=''):
    """
    Present final extracted data on a slide

    @param data: @2DArray with first entry as the column titles and rest of the entries as the subsequent rows.
    @param sz: number of rows to be displayed per slide.
    """
    n = len(data)

    # dev mode only
    if limit != -1:
        n = min(limit, n)
    
    # break data into chunks of size @param sz per slide
    for i in range(1, n, sz):
        # data[0] are the column titles
        entries = [data[0]]
        # break the chunk
        entries += data[i:i+sz]
        naming = title
        if i > 1: 
            naming += ' Continued...'
        # create the PPT slide
        create_a_slide_with_data(prs, entries, titleofslide=naming, global_null_value=global_null_value)

def create_a_multiselect(headers, key, title=''):
    """
    Create a multiselect option

    @param headers: The choices to be displayed in the multiselect option.
    @param key: [Required for Streamlit API], unique name of the multiselect widget.
    @param title: [Required for Streamlit API], Title to be displayed over multiselect widget.
    """

    opencase_container = st.container()
    # Checkbox to select all the columns at once.
    want_all_open = st.checkbox("Select all", key=key)
    if want_all_open:
        opcsheads = opencase_container.multiselect(
            title,
            headers, headers, key=key+'_ms', help='Select the columns to be added in the table'
        )
    else:
        # Select columns one by one
        opcsheads =  opencase_container.multiselect(
            title,
            headers, key=key+'_ms', help = 'Select the columns to be added in the table'
        )
    # return selected columns and the trigger associated with the checkbox
    return opcsheads, want_all_open

def end_slide():
    """
    Append the end slide to the PPT

    """
    layout = prs.slide_layouts[5]
    slide=prs.slides.add_slide(layout)
    # slide.shapes.add_picture('data/end.png', 0, 0, prs.slide_width, prs.slide_height)


def commence_ppt_creation():
    """
    Generates PPT slides for Charts, Closed cases, Open cases and Extra slides as per the requirement.

    """
    # extract data if the column headers are valid
    if len(opcsheads) > 0 or len(cscsheads) > 0:
        opencases, closedcases = extract_rows(dataframe, True)
    if len(opcsheads) == 0:
        opencases = []
    if len(cscsheads) == 0:
        closedcases = []
    
    # display charts onto the slides
    for ele in ext_cha_arr:
        data_col = ele[0]
        chart_type = ele[1].lower()
        title = ele[2]
        # get the data to be displayed onto the chart
        data = convert_to_categories(dataframe[data_col], global_null_value,True)
        if chart_type == 'pie':
            create_a_chart(prs, data, [2.5,1,8,6], None, typeOfChart=chart_type, nameofchart=title)
        elif chart_type == 'bar':
            create_a_chart(prs, data, [0.5,0.6,12,6.5], None, typeOfChart=chart_type, nameofchart=title)
        elif chart_type == 'donut':
            create_a_chart(prs, data, [0.5,0.6,12,6.5], None, typeOfChart=chart_type, nameofchart=title)

    # create open cases and closed cases slide if valid    
    if not (openCaseCol is None or openCaseCol == '<select>'):
        present_on_slide(opencases, openCase_rows_per_slide, 'Open Cases')
        present_on_slide(closedcases, closeCase_rows_per_slide, 'Closed Cases')
    
    # Create extra tables that were selected in the UI
    for ele in ext_tab_arr:
        data = ele[0]
        data = extract_rows(dataframe, False, data)
        title = ele[1]
        rows_per_slide = ele[2]
        present_on_slide(data, rows_per_slide, title)
    
    end_slide()

    # Save the PPT 
    prs.save("x_final_ppt.pptx")
    
    # download the PPT
    with open("x_final_ppt.pptx", "rb") as file:
        btn = st.download_button(
            label="Download PPT",
            data=file,
            file_name="pyStAR_final.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

def remove_keys(keys):
    """
    Removes keys from the session_state.

    @param keys: @Array containing keys that needs to be deleted from st.session_state
    """
    for key in keys:
        if key in st.session_state:
            del st.session_state[key]

if __name__ == "__main__":
    # get the uploaded file
    st.markdown("## Excel File")
    uploaded_file = st.file_uploader("Upload your file", type=['xlsx','xlsm'], accept_multiple_files=False, key='upl', on_change=lambda:remove_keys(['dataframe','headers']))
    st.markdown("## PPT Template")
    template_chk = st.checkbox("Do you want to upload a custom template file?", key='temp_chk', help="If you don't want to upload a template, a default template will be used.")
    if template_chk:
        st.info('A template basically makes up the initial slides in your final PPT. All the slides will be appended to this template.')
        template_file = st.file_uploader("Upload template PPT file", type=['pptx'], accept_multiple_files=False, key='temp_upl')
    if template_file:
        prs = Presentation(template_file)
    else:
        prs = Presentation(default_template)

    # if uploaded file is present proceed further with the UI
    if uploaded_file:
        st.markdown("## Data fetching 📊")
        
        with st.expander('Want to change default null value replacement for the data?' ):
            global_null_value = st.text_input("Type the replacement value.", value='-', help="[OPTIONAL] All the null values in excel file will be replaced by this character. By Default, they are being replaced by ' -   '. ")

        user_start_location = st.text_input('Enter Start Cell', placeholder='A13', key='user_start_loc',on_change=lambda:remove_keys(['dataframe','headers']),  help="Enter the cell address from where you want the data to be extracted from the excel file.")
        if user_start_location != "":   
            
            # IMPORTANT: Before we understand what session_state is, we need to understand that streamlit
            # refreshes the entire UI everytime a change is detected in UI elements. So in order to save 
            # unnecessary recomputation of time expensive events, we can save them as session variables in
            # session_state dictionary

            # since data fetching is expensive, we can fetch it once and store it as a session variable
            # now, everytime the UI refreshes itself the session_state variables are reused instead of recalculation.

            # Fetching data and storing the dataframe as a session variable
            if 'dataframe' not in st.session_state:
                # if dataframe is not in session variables, calculate it and store it
                cols, rows = split_start_address(user_start_location)    
                beginx = int(rows)
                beginy = convert_excel_col_number(cols)
                dataframe = get_data(uploaded_file)
                dataframe = dataframe.loc[:, dataframe.columns.notna()]
                dataframe = dataframe.reset_index(drop=True)
                headers = dataframe.columns.tolist()
                st.session_state['dataframe'] = dataframe
                st.session_state['headers'] = headers
            else:
                # if dataframe is available in session variables, reuse it instead of recalculation
                dataframe = st.session_state['dataframe']
                headers = st.session_state['headers']
            
            

            st.markdown("## Data Preview 📋")
            
            display_head = dataframe.head()
            # Error prone area, The starting address maybe invalid or may contain location that is not valid 
            # or may contain duplicate columns, so handled them here.
            try:
                st.dataframe(display_head)
                st.info("Please verify if all the columns titles are displayed in the preview below, if not please re-enter correct starting cell address.")
            except Exception as e:
                # assuming error occured while calculating preview, maybe due to duplicate columns
                _dict = {}
                # verify if there are any duplicate column names
                flag_duplicate = False
                for ele in display_head:
                    if ele in _dict: 
                        _dict[ele] += 1
                    else:
                        _dict[ele] = 1
                dup_col = []
                _display_head = []
                for key,value in _dict.items():
                    if value > 1:
                        flag_duplicate = True
                        dup_col.append([key, value])
                    _display_head.append(key)
                
                if flag_duplicate:
                    st.error("ERROR: There maybe something wrong with the column names, possibly duplicate column names, please enter correct starting cell address or review your data")
                    show_this = f"The following column(s) are present more than once: "
                    for ele in dup_col:
                        show_this += f"{ele[0]} exists {ele[1]} times, "
                    show_this = show_this[:-2]
                    st.error(show_this)
                else:
                    _display_head = dataframe[_display_head]
                    st.dataframe(_display_head)
                    st.info("Please verify if all the columns titles are displayed in the preview below, if not please re-enter correct starting cell address.")
            
            
            st.markdown("## Create Slides 📰")
            with st.expander("Do you want to split the slides into open cases and closed cases?"):
                st.markdown("### Case Split")
                
                openCaseCol = st.selectbox('Select the case split column: ', ['<select>'] + headers, help="Enter a column title on the basis of which the segregation will be performed. Usually it is either “Open Days” column or “Open Cases” column")
                if openCaseCol == '<select>':
                    st.write("Please Select an open case column.")
                else:
                    global_compare_false_val = st.text_input("Enter the seperation delimiter [Case Sensitive]", value='-', help="The slides will be split based on this delimiter, if entry in a cell equals to this delimiter, it goes onto the closed case slide. It can also be a list of commas seperated values: -,NO,_.").split(',')

                    opcsheads, open_case_trigger = create_a_multiselect(headers, "open_cases", "Select column for open case slide: ")
                    op_chk = st.checkbox("Do you want to insert blank columns?", key='op_chkbox', help="Enter the new column title in the space provided below. You can also add multiple columns by separating the titles with commas. For e.g. Col A, Col B.")
                    if op_chk:
                        opcsheads = st.text_area("", ",".join(opcsheads)).split(',')

                    openCase_rows_per_slide = st.number_input("Enter number of rows per slide: ", 0, 1000, 6, 1, key="op_per_slide", help="Enter number of rows to be displayed in each slide.")

                    cscsheads, close_case_trigger = create_a_multiselect(headers, "close_cases", "Select column for close case slide: ")
                    
                    cs_chk = st.checkbox("Do you want to insert blank columns?", key='cs_chkbox',help="Enter the new column title in the space provided below. You can also add multiple columns by separating the titles with commas. For e.g. Col A, Col B.")
                    if cs_chk:
                        cscsheads = st.text_area("", ",".join(cscsheads)).split(',')

                    closeCase_rows_per_slide = st.number_input("Enter number of rows per slide: ", 0, 1000, 6, 1, key="cs_per_slide", help="Enter number of rows to be displayed in each slide.")
                          

            with st.expander("Do you want to add more tables to the PPT slides?"):
                st.markdown("### Add the tables")
                ext_num = st.number_input("Enter the number of extra tables", 0, 1000, 0, 1, key="ext_num")
                
                # Sometimes the input may return a float value, so we convert it into an int
                try:
                    ext_num = int(ext_num)
                except:
                    ext_num = 0

                # Create ext_num amount of UI widgets for taking inputs
                ext_tab_arr = []
                for i in range(ext_num):

                    tab_name = st.text_input("Enter title of the slide",value="Table "+str(i+1), key="tab_name_"+str(i+1))
                    temp, temp_trigger = create_a_multiselect(headers, "tab_"+str(i+1), "Select columns for this table:")

                    ext_chk = st.checkbox("Do you want to insert blank columns?", key='ext_chkbox'+str(i+1))
                    
                    if ext_chk:
                        temp = st.text_area("", ",".join(temp), key="ext_chkbox_"+str(i+1)).split(',')
                    
                    sz_per_slide = st.number_input("Enter number of rows per slide: ", 0, 1000, 6, 1, key="sz_per_slide"+str(i+1))

                    if len(temp) == 0:
                        st.error("Please select at least 1 column.")
                    else:
                        ext_tab_arr.append([temp, tab_name, sz_per_slide])

            # Charts
            with st.expander("Do you want to add charts?"):
                st.markdown("### Charts")
                ext_num_c = st.number_input("Enter the number of charts you want", 0, 1000, 0, 1, key="ext_num_c")
                try:
                    ext_num_c = int(ext_num_c)
                except:
                    ext_num_c = 0
                choices_of_charts = ['Pie', 'Bar','Donut']
                ext_cha_arr = []
                for i in range(ext_num_c):
                    cha_name = st.text_input("Enter chart name",value="Chart "+str(i+1), key="cha_name_"+str(i+1))
                    temp = st.selectbox("Select column for this chart:", headers, key="cha_"+str(i+1), help='The chart will be created based on data of this column.')
                    type_of_chart = st.selectbox('Choose the type of chart.', choices_of_charts, key="cha_choices_"+str(i+1))
                    ext_cha_arr.append([temp, type_of_chart,cha_name])

            st.markdown("## Submit and create PPT ✔️")
            # Button to submit UI data for PPT creation
            butt_trigger = st.button('Submit')
            if butt_trigger:
                commence_ppt_creation()

components.html("<!-- HVLC was here -->")
