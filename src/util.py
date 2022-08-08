from requests import head
#for pthon3.10
import collections 
import collections.abc

# for pptx
from pptx.util import Inches, Cm, Pt
from pptx.dml.color import ColorFormat, RGBColor
from pptx.enum.text import MSO_ANCHOR
from pptx.oxml.xmlchemy import OxmlElement

# for charts
from pptx.chart.data import CategoryChartData, ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION, XL_TICK_MARK,XL_TICK_LABEL_POSITION
from pptx.enum.text import PP_ALIGN

def split_start_address(user_start):
    """
    Split Start address into row-number and col-number

    @param user_start: e.g. A13, Start address of the dataframe
    @return: cols: A, rows: 13
    """
    user_start = user_start.upper()
    n = len(user_start)
    num_ind = -1
    for i in range(n):
        idx = ord(user_start[i])
        if idx >= ord('0') and idx <= ord('9'):
            num_ind = i
            break
    cols = user_start[:num_ind]
    rows = user_start[num_ind:]
    return cols, rows

def convert_excel_col_number(current_letter):
    """
    Utility function to convert Column name into a number: e.g. AA->27

    @param current_letter: Column name as letters of Alphabet, e.g.: AAX
    @return: Numeric equivalent of the title.
    """
    n = len(current_letter)
    mul = 1
    sm = 0
    for i in range(n-1, -1, -1):
        sm += mul*(ord(current_letter[i]) - ord('A') + 1)
        mul *= 26
    return sm

def convert_number_excel_col(sm):
    """
    Utility function to convert a column number into a : e.g. AA->27

    @param current_letter: Column name as letters of Alphabet, e.g.: AAX
    @return: Numeric equivalent of the title.
    """
    new_col = ""
    while sm > 0:
        sm-=1
        new_col += chr(ord('A') + sm%26)
        sm = sm//26
    new_col = new_col[::-1]
    return new_col

def add_to_column_letter(current_letter, add_val):
    """
    Add @param add_val to current letter(column address).
    e.g. AA + 5 : AF

    @param current_letter: 
    @return: updated column address as a string
    """
    # forward:  A-Z: 1:26
    # backward: A-Z: 0:25

    # convert to number
    sm = add_val + convert_excel_col_number(current_letter)
    # convert to string
    new_col = convert_number_excel_col(sm)
    return new_col

def comparator_break(curr_val, true_val):
    """
    Compare @param curr_val with elements in @param true_val

    @param curr_val: Value to be compared 
    @param true_val: @Array containing true values against which @param curr_val is compared.
    @return: boolean value
    """
    for ele in true_val:
        if curr_val == ele:
            return True
    return False


def add_logo(prs, slide):
    """
    Utility function to add a logo to @param prs presentation in @param slide slide.

    @param prs: PPT presentation object.
    @param slide: PPT slide object on which logo is to be displayed.
    """

    slide_idx = prs.slides.index(slide)
    shapes = prs.slides[slide_idx].shapes
    logo = shapes.add_picture('data/logo.jpg', Cm(32.8), Cm(0.2), Cm(0.72), Cm(1))


def SubElement(parent, tagname, **kwargs):
    """
    

    @param file: Path to the excel file
    @return: Pandas dataframe
    """
    element = OxmlElement(tagname)
    element.attrib.update(kwargs)
    parent.append(element)
    return element

def _set_cell_border(cell, border_color="000000", border_width='12700'):
    """
    

    @param file: Path to the excel file
    @return: Pandas dataframe
    """
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for lines in ['a:lnB']:
        ln = SubElement(tcPr, lines, w=border_width, cap='flat', cmpd='sng', algn='ctr')
        solidFill = SubElement(ln, 'a:solidFill')
        srgbClr = SubElement(solidFill, 'a:srgbClr', val=border_color)
        prstDash = SubElement(ln, 'a:prstDash', val='solid')
        round_ = SubElement(ln, 'a:round')
        headEnd = SubElement(ln, 'a:headEnd', type='none', w='med', len='med')
        tailEnd = SubElement(ln, 'a:tailEnd', type='none', w='med', len='med')



def create_a_chart(prs, data, position, slide, typeOfChart, nameofchart=''):
    """
    Creates a chart on @param slide else add a new slide.

    @param prs: PPT Presentation object
    @param data: @2DArray, first entry is for categories, second entry is the respective values. e.g. [x-axis_data, y-axis_data]
    @param position: @Array containing 4 elements: [left, top, width, height] in inches.
    @param slide: slide object, None if slide is to be appended to existing presentation else pass the slide object.
    @param typeofchart: type of chart, e.g: bar,pie,donut.
    @param nameofchart: Title of Chart that appears on PPT slide.
    """

    # Blank slide with title, is susceptible to change(depends on layout of the template).
    layout = prs.slide_layouts[6]
    if slide is None:
        slide=prs.slides.add_slide(layout)

    
    add_logo(prs, slide)
    
    # initialize Chart object
    chart_data = ChartData()
    chartType = XL_CHART_TYPE.PIE

    # x: left, y: top, cx: width, cy: height
    x, y, cx, cy = Inches(position[0]), Inches(position[1]), Inches(position[2]), Inches(position[3])

    # Add chart categories and values
    chart_data.categories = data[0]
    chart_data.add_series(nameofchart, tuple(data[1]))

    # identify chart type
    if typeOfChart == 'bar':
        # Establish Chart type
        chartType = XL_CHART_TYPE.COLUMN_CLUSTERED
        
        # create the chart
        chart = slide.shapes.add_chart(
            chartType, x, y, cx, cy, chart_data
        ).chart

        # Cosmetic changes to Chart
        yaxis = chart.value_axis
        xaxis = chart.category_axis
        # Tick marks
        yaxis.tick_label_position = XL_TICK_LABEL_POSITION.NONE
        yaxis.major_tick_mark = XL_TICK_MARK.NONE
        yaxis.has_major_gridlines = False
        xaxis.has_minor_gridlines = False
        # tick_labels = yaxis.tick_labels
        # tick_labels.number_format = '0.0%'
        # tick_labels.font.size = Pt(13)
        xlabels = xaxis.tick_labels
        xlabels.font.size = Pt(11)

        # chart legend
        # print(tick_labels.number_format)
        # chart.has_legend = True
        # chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        # chart.legend.include_in_layout = False
        
        chart.plots[0].has_data_labels = True
        data_labels = chart.plots[0].data_labels
        data_labels.number_format = '0.0%'
        data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
        data_labels.font.size = Pt(11)
    elif typeOfChart == 'pie':
        # Establish chart type
        chartType = XL_CHART_TYPE.PIE

        # Create the chart
        chart = slide.shapes.add_chart(
            chartType, x, y, cx, cy, chart_data
        ).chart

        # Cosmetic changes to the chart
        # chart.chart_style = 12

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        
        chart.plots[0].has_data_labels = False
        # data_labels = chart.plots[0].data_labels
        # data_labels.number_format = '0.0%'
        # data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
        # data_labels.font.size = Pt(12)
    elif typeOfChart=='donut':
        # Establish chart type
        chartType = XL_CHART_TYPE.DOUGHNUT

        # Create the chart
        chart = slide.shapes.add_chart(
            chartType, x, y, cx, cy, chart_data
        ).chart

        # Cosmetic changes to the chart
        # chart.chart_style = 12

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        
        chart.plots[0].has_data_labels = False
        # data_labels = chart.plots[0].data_labels
        # data_labels.number_format = '0.0%'
        # data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
        # data_labels.font.size = Pt(12)
    return

def create_a_slide_with_data(prs, data, titleofslide='',global_null_value='-',sizes=None):
    """
    Create a slide and present the data.

    @param prs: PPT Presentation object
    @param data: @2DArray where first entry is the column headers and rest of the entries are the rows to be displayed on a slide.
    @param titleofslide: Title of the slide to be displayed in the PPT.
    @param sizes: Experimental feature for variable size columns on the PPT slide.[Not Implemented]
    """
    # n: Total number of rows, m: Total number of columns
    n = len(data)
    m = len(data[0])

    # Blank slide layout with title, is susceptible to change(depends on layout of the template).
    layout = prs.slide_layouts[5]
    slide=prs.slides.add_slide(layout)
    add_logo(prs, slide)

    # Slide Cosmetics
    title=slide.shapes.title
    title.text=titleofslide
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    title.text_frame.paragraphs[0].font.size = Pt(32)
    # Red color RGB
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF,0x00,0x00)
    title.text_frame.paragraphs[0].font.name = 'Adobe Clean'
    # title position
    # title.top = Cm(1)
    title.left = Cm(1)
    title.width = Inches(10)
    title.height = Inches(1)

    # x: left, y: top, cx: width, cy: height
    x, y, cx, cy = Inches(0.3), Inches(1.25), Inches(12.8), Inches(0.5)
    table = slide.shapes.add_table(n, m, x, y, cx, cy).table
    
    # Experimental feature for variable column width
    # for i in range(len(sizes)):
    #     table.columns[i].width = Inches(sizes[i])

    # For every cell in the table being displayed, map data to corresponding cell.
    for i in range(n):
        for j in range(m):
            cell = table.cell(i, j)
            value = str(data[i][j])

            # If value of data is nan, replace it with global_null_value
            if value == 'nan':
                value = global_null_value
            
            # Cosmetic changes for table cells.
            cell.text = value
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(10.5)
            paragraph.alignment = PP_ALIGN.CENTER
            paragraph.font.color.rgb = RGBColor(0x00,0x00,0x00)
            paragraph.font.name = 'Adobe Clean'
            if i != 0:
                _set_cell_border(cell, "e2e2e2")
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xE2, 0xE2, 0xE2)


def convert_to_categories(dataframe, global_null_value='-', div=False):
    """
    Break data into key,value pairs where keys are the categories and value is the frequency of that particular category in data.
    
    @param dataframe: Dataframe object.
    @param global_null_value: Global Null value.
    @param div: Normalize the frequencies to fit a range(100% here).
    """

    # Make frequency dictionary of data, where key is a category and value is its frequency.
    _dict ={}
    for name,value in dataframe.iteritems():
        if str(value) == global_null_value:
            continue
        if value in _dict:
            _dict[value]+=1
        else:
            _dict[value]=1
    # Break the frequency dictionary into a 2D array where first entry is the categories and second entry is the frequencies associated with it.
    data = [[],[]]
    for key,value in _dict.items():
        data[0].append(key)
        data[1].append(value)
    if div:
        sm = sum(data[1])
        for i in range(len(data[1])):
            data[1][i] /= sm
    return data

"""
HVLC was here
"""